"""Generate Week-8 progress slides for the blockchain-consortium-networks SLR.

Mirrors the structure of slides/week7-progress-2026-04-14.pptx but with the
agent full-text review pipeline, the final PRISMA funnel (279 include / 584
exclude across 863 candidates), the exclude-reason breakdown, and the
data-extraction/quality-assessment plan for Week 9.
"""

from __future__ import annotations

import csv
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SCREEN_CSV = ROOT / "data" / "processed" / "fulltext_screening.csv"
OUTPUT = ROOT / "slides" / "weekly-progress-2026-04-21.pptx"

ACCENT = RGBColor(0x1F, 0x3A, 0x68)        # navy
ACCENT_LIGHT = RGBColor(0x4C, 0x78, 0xA8)  # mid blue
GOOD = RGBColor(0x2E, 0x7D, 0x32)          # green
WARN = RGBColor(0xE6, 0x8A, 0x00)          # orange
BAD = RGBColor(0xC6, 0x28, 0x28)           # red
GREY = RGBColor(0x55, 0x55, 0x55)


def load_counts() -> dict:
    rows = list(csv.DictReader(SCREEN_CSV.open(encoding="utf-8")))
    total = len(rows)
    ta = Counter(r["title_abstract_decision"] for r in rows)
    ftd = Counter(r["fulltext_decision"] for r in rows)
    final = Counter(r["final_decision"] for r in rows)
    excl_reason = Counter(
        r["fulltext_exclusion_reason_code"]
        for r in rows
        if r["final_decision"] == "exclude"
    )
    return {
        "candidates": total,
        "ta_include": ta.get("INCLUDE", 0),
        "ta_uncertain": ta.get("UNCERTAIN", 0),
        "ft_include": ftd.get("include", 0),
        "ft_exclude": ftd.get("exclude", 0),
        "final_include": final.get("include", 0),
        "final_exclude": final.get("exclude", 0),
        "excl_reasons": excl_reason,
    }


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT
    bg.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.7), Inches(2.2),
                                  Inches(12), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Week 8 Progress Report"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf.add_paragraph()
    p2.text = "Blockchain in Consortium Networks - Systematic Literature Review"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xCF, 0xDD, 0xEE)

    p3 = tf.add_paragraph()
    p3.text = "Tony Eneh - Apr 21, 2026"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0xCF, 0xDD, 0xEE)


def add_section_title(slide, text: str) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide.part.package.presentation_part
        .presentation.slide_width, Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.1),
                                  Inches(12), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_bullet_slide(prs: Presentation, title: str, bullets: list) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, title)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.95),
                                  Inches(12.3), Inches(6.3))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            level, txt = item
        else:
            level, txt = 0, item
        p.text = txt
        p.level = level
        p.font.size = Pt(20 if level == 0 else 16)
        p.font.color.rgb = ACCENT if level == 0 else GREY
        p.space_after = Pt(6)


def add_metric_card(slide, left, top, width, height, label, value, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.text = str(value)
    p1.alignment = 2  # center
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = tf.add_paragraph()
    p2.text = label
    p2.alignment = 2
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_funnel_slide(prs: Presentation, c: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "PRISMA 2020 - Final Screening Funnel")

    cards = [
        ("Records identified", "1,707", ACCENT_LIGHT),
        ("After dedup", "1,638", ACCENT_LIGHT),
        ("Title/abstract candidates", str(c["candidates"]), ACCENT),
        ("Final included", str(c["final_include"]), GOOD),
        ("Final excluded", str(c["final_exclude"]), BAD),
    ]
    width = Inches(2.3)
    height = Inches(1.5)
    top = Inches(1.2)
    left = Inches(0.5)
    for i, (label, value, color) in enumerate(cards):
        add_metric_card(slide, left + i * (width + Inches(0.15)),
                        top, width, height, label, value, color)

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(3.0),
                                  Inches(12.3), Inches(4))
    tf = tx.text_frame
    tf.word_wrap = True
    items = [
        f"863 candidates = {c['ta_include']} first-pass INCLUDE + "
        f"{c['ta_uncertain']} UNCERTAIN (resolved by agent pipeline this week).",
        f"Agent full-text review of all 863: include = {c['ft_include']}, "
        f"exclude = {c['ft_exclude']}.",
        "Pipeline: OA discovery (arXiv -> OpenAlex -> raw URL) -> in-memory "
        "PDF/HTML extraction (pypdf + bs4) -> heuristic scoring "
        "(institutional + implementation + evidence terms) -> conservative "
        "finalisation (needs_human_check -> exclude).",
        "Single source of truth: data/processed/fulltext_screening.csv.",
        "Auto-generated TikZ flow figure: paper/figures/prisma_flow.tex.",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = GREY
        p.space_after = Pt(6)


def add_exclusion_chart_slide(prs: Presentation, c: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Why we excluded 584 candidates")

    label_map = {
        "insufficient_fulltext_access": "Insufficient full-text access",
        "no_full_text_access": "No open-access full text",
        "second_pass_likely_exclude": "Second-pass likely exclude",
        "agent_uncertain_full_text": "Agent borderline (conservative)",
        "no_quantitative_or_reproducible_evidence": "No reproducible evidence",
        "policy_or_commentary_only": "Policy / commentary only",
        "borderline_fulltext_case": "Borderline full-text case",
        "not_institutional_setting": "Not institutional setting",
        "no_implementation_detail": "No implementation detail",
    }
    items = c["excl_reasons"].most_common()
    cats = [label_map.get(k, k or "(unspecified)") for k, _ in items]
    vals = [v for _, v in items]

    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series("Excluded", vals)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5), Inches(1.0), Inches(8.5), Inches(6.0),
        chart_data,
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(11)

    tx = slide.shapes.add_textbox(Inches(9.2), Inches(1.0),
                                  Inches(4.0), Inches(6.0))
    tf = tx.text_frame
    tf.word_wrap = True
    bullets = [
        ("Top driver", True),
        ("Open-access gap dominates "
         "(491 of 584; 84%).", False),
        ("Author work this week", True),
        ("Procure full-text for the 491 "
         "paywalled rows via institutional "
         "access; re-run agent.", False),
        ("Spot-check the 32 borderline "
         "agent exclusions and the 48 "
         "second-pass excludes.", False),
    ]
    for i, (txt, head) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(14) if head else Pt(12)
        p.font.bold = head
        p.font.color.rgb = ACCENT if head else GREY
        p.space_after = Pt(4)


def main() -> None:
    c = load_counts()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_bullet_slide(prs, "Week 8 Goals", [
        "Close the title/abstract uncertain queue end-to-end (no manual backlog).",
        ("    Second-pass triage -> OA discovery -> agent review -> "
         "finalisation."),
        "Run agent full-text review on the 571 first-pass INCLUDE records.",
        "Generate the PRISMA 2020 flow diagram from data, not by hand.",
        "Refresh paper draft and TODO with Week-8 numbers.",
        "Lock the candidate pool so Week 9 can start data extraction.",
    ])

    add_funnel_slide(prs, c)

    add_bullet_slide(prs, "Pipeline Built This Week (all reproducible)", [
        "scripts/triage_uncertain_records.py - second-pass NLP triage of 292 uncertain rows.",
        "scripts/split_uncertain_triage.py - partitions into likely_include / likely_exclude / needs_manual.",
        "scripts/discover_ambiguous_fulltexts.py - arXiv + OpenAlex OA discovery (15 PDFs, 38 links, 68 paywalled).",
        "scripts/prioritize_ambiguous_fulltexts.py - ranks the 121 needs-manual records by access bucket.",
        "scripts/review_accessible_fulltexts.py - agent review of 53 accessible papers (20/1/32).",
        "scripts/finalize_uncertain_decisions.py - rule-based promotion: 143 include / 149 exclude.",
        "scripts/review_firstpass_fulltexts.py - NEW. Agent review of 571 first-pass includes (136/435), with checkpoint+resume.",
        "scripts/generate_prisma_flow.py - NEW. Emits paper/figures/prisma_flow.tex from CSV counts.",
    ])

    add_exclusion_chart_slide(prs, c)

    add_bullet_slide(prs, "Paper & Repo Updates", [
        "paper/main.tex: tikz packages added; abstract + Study Selection + Study Characteristics rewritten with Week-8 numbers.",
        "paper/figures/prisma_flow.tex: auto-generated TikZ figure (1707 -> 1638 -> 863 -> 279 / 584).",
        "TODO.md: 714-record stale pool replaced; Immediate Next Actions repointed at Step 6 (extraction) and Step 7 (quality).",
        "Atomic commits this week (3 + previous 7 from week 7 batch):",
        "    feat(screen): agent full-text review of 571 first-pass title/abstract includes",
        "    feat(paper): auto-generated PRISMA 2020 flow diagram (TikZ)",
        "    docs(slr): refresh TODO and paper with Week-8 final screening metrics",
        "All Python scripts deterministic; rerunning the pipeline reproduces the exact 279/584 split.",
    ])

    add_bullet_slide(prs, "Risks & Limitations", [
        "Agent-only screening is the largest risk vector.",
        ("    491 of 584 exclusions are driven by lack of OA full text - "
         "could hide eligible studies behind paywalls."),
        ("    32 borderline cases were conservatively excluded; sampling needed."),
        "Heuristic scoring uses keyword baskets, not full LLM reasoning.",
        ("    Acceptable for first-pass agent screening; author confirmation "
         "still required before final inclusion in synthesis."),
        "Snowball sampling on the 279 includes is not yet started.",
    ])

    add_bullet_slide(prs, "Week 9 Plan", [
        "Step 6 - Data extraction on the 279 included studies.",
        ("    Populate data/processed/data_extraction.csv (network type, "
         "consensus, governance, deployment scale, evidence type)."),
        "Step 7 - Quality assessment.",
        ("    Apply CASP-style rubric; flag low-quality items for "
         "sensitivity analysis."),
        "Snowball-sample references of the 279 includes.",
        "Procure paywalled full-texts via institutional access; re-run agent reviewer to recover false negatives.",
        "Begin taxonomy synthesis draft (Discussion section scaffolding).",
    ])

    add_bullet_slide(prs, "Asks", [
        "Confirm the conservative finalisation policy (needs_human_check -> exclude) before lock.",
        "Approve the move to Step 6 (data extraction) over the 279-record subset.",
        "Recommend any priority venues / authors for snowball sampling.",
        "Confirm CASP-style quality rubric is acceptable, or suggest an alternative.",
    ])

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")
    print(f"  candidates       : {c['candidates']}")
    print(f"  final include    : {c['final_include']}")
    print(f"  final exclude    : {c['final_exclude']}")
    print(f"  excl reasons     : {dict(c['excl_reasons'])}")


if __name__ == "__main__":
    main()
